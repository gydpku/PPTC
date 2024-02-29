import collections 
import collections.abc
import pptx.parts.image
import pptx.enum.shapes as shapes
from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR, MSO_FILL
from pptx.dml.color import RGBColor
from .content_selection import *
from .openai_api import *
from .prompt_factor import *

global slides
global ppt
global global_args
SCALE = 1000
slide_height = 0
slide_width = 0
shape_list = ['PLACEHOLDER', 'PICTURE', 'CHART', 'TABLE', 'TEXT_BOX', 'AUTO_SHAPE']
slides = None


def get_fill_color(shape):
    if shape.fill.type == 1:  # Solid fill
        color = shape.fill.fore_color
        if hasattr(color, "rgb"):
            return color.rgb
    return None

def set_slides(slides1):
    global slides
    slides = slides1

class BasicShape:
    def __init__(self, shape):
        self.shape_type = shape.shape_type
        self.height = shape.height // SCALE
        self.width = shape.width // SCALE
        self.left = shape.left // SCALE
        self.top = shape.top // SCALE
        self.name = shape.name
        self.shape_id = shape.shape_id
    
    @property
    def text_info(self):
        pass

    @property
    def space_info(self):
        return f"Visual Positions: left={self.left}, top={self.top}\n"
    
    @property
    def size_info(self):
        return f"Size: height={self.height}, width={self.width}\n"

    @property
    def style_info(self):
        pass

    @property
    def discription(self):
        # return f"[{self.name.split(' ')[0]}]\n" 
        try:
            assert not self.name.split(' ')[0] == "Google"
            return f"[{self.name.split(' ')[0]}]\n"
        except:
            return f"[{str(self.shape_type).split(' ')[0].strip()}]\n" 

    def __repr__(self):
        s = ""
        s += self.discription
        s += self.size_info
        if self.text_info is not None:
            s += self.text_info
        s += self.space_info
        if self.style_info is not None:
            s += self.style_info
        return s


class Picture(BasicShape):
    def __init__(self, shape, id=None):
        super().__init__(shape)
        self.image = shape.image
        self.rotation = int(shape.rotation)
        self.id = id
    
    @property
    def style_info(self):
        return f"Picture Style: rotation={self.rotation}\n"
    
    @property
    def discription(self):
        if self.id != None:
            return f"[Picture {self.id}]\n"
        else:
            return f"[Picture]\n"

class Table(BasicShape):
    def __init__(self, shape):
        super().__init__(shape)
        self.table = shape.table
        self.rows = shape.table.rows
        self.columns = shape.table.columns

    @property
    def text_info(self):
        s = "Data:\n"
        for row in self.rows:
            s += "|"
            for col in row.cells:
                s += f"{col.text}|"
            s += "\n"
        return s
    
    @property
    def discription(self):
        return f"[Table] with {len(self.rows)} rows and {len(self.columns)} columns\n" 

class Chart(BasicShape):
    def __init__(self, shape):
        super().__init__(shape)
        self.chart = shape.chart
        self.title = shape.chart.chart_title.text_frame.text
        self.chart_type = str(shape.chart.chart_type).split(' ')[0]
    
    @property
    def text_info(self):
        s = ""
        if self.title:
            s += f"Title: {self.title}\n"
        s += f"Chart Type: {self.chart_type}\n"
        s += "Data:\n"
        try:
            for series in self.chart.series:
                s += f"{series.name}: "
                for value in series.values:
                    s += f"{value}, "
                s += "\n"
        except:
            pass
        return s
    
    @property
    def style_info(self):
        return ""
    
    @property
    def discription(self):
        return "[Chart]\n"


class Textbox(BasicShape):
    def __init__(self, shape, id=None):
        super().__init__(shape)
        self.text = shape.text_frame.text
        self.paragraphs = shape.text_frame.paragraphs
        try:
            self.font = self.paragraphs[0].runs[0].font
        except:
            self.font = self.paragraphs[0].font
        self.bold = self.font.bold
        self.italic = self.font.italic
        self.underline = self.font.underline
        self.size = self.font.size if self.font.size is not None else self.paragraphs[0].font.size
        try:
            self.color = self.font.color.rgb 
        except:
            self.color = None
        self.fill = get_fill_color(shape)
        self.font_name = self.font.name
        self.line_spacing = self.paragraphs[0].line_spacing
        self.align = self.paragraphs[0].alignment
        self.id=id
    
    @property
    def text_info(self):
        return f"Text: {self.text}\n"
    
    @property
    def style_info(self):
        return f'Font Style: bold={self.bold}, italic={self.italic}, underline={self.underline}, size={self.size}, color={self.color}, fill={self.fill}, font style={self.font_name}, line_space={self.line_spacing}, align={self.align}\n'

    @property
    def discription(self):
        if self.id != None:
            return f"[TextBox {self.id}]\n"
        else:
            return f"[TextBox]\n"
    
class Placeholder(BasicShape):
    def __init__(self, shape):
        super().__init__(shape)
        self.fill = get_fill_color(shape)
        self.text = shape.text_frame.text
        if shape.has_text_frame:
            textframe = shape.text_frame
            try:
                font = shape.text_frame.paragraphs[0].runs[0].font
            except:
                font = shape.text_frame.paragraphs[0].font
            self.bold = font.bold
            self.italic = font.italic
            self.underline = font.underline
            self.size = font.size 
            try:
                self.color = font.color.rgb 
            except:
                self.color = None
            self.font_name = font.name
            self.line_spacing = textframe.paragraphs[0].line_spacing
            self.align = textframe.paragraphs[0].alignment
    
    @property
    def text_info(self):
        if self.text is not None:
            return f"Text: \n{self.text}\n"
        else:
            return ""
    
    @property
    def style_info(self):
        return f'Font Style: bold={self.bold}, italic={self.italic}, underline={self.underline}, size={self.size}, color={self.color}, fill={self.fill}, font style={self.font_name}, line_space={self.line_spacing}, align={self.align}\n'


class AutoShape(BasicShape):
    def __init__(self, shape):
        super().__init__(shape)
        self.text = shape.text_frame.text
        self.fill = get_fill_color(shape)
    
    @property
    def text_info(self):
        return f"Text: {self.text}\n"
    
    @property
    def style_info(self):
        return f"Shape Style: fill={self.fill}\n"
        # return ""

def hasshape(shape_str, shape_list):
    for shape in shape_list:
        if shape in shape_str:
            return True
    return False

def get_content(need_text,need_style,need_position,need_title,need_content,need_picture,need_table,need_chart,need_textbox,need_shape):
    global slides
    global global_args
    s = ""
    idx = 0
    for slide in slides:
        s += f"Slide {idx} with background color {get_fill_color(slide.background)}:\n"
        if slide.notes_slide.notes_text_frame.text:
            s += f"Notes: {slide.notes_slide.notes_text_frame.text}\n"
        idx += 1
        textbox_idx = 0
        picture_idx = 0

        for shape in slide.shapes:
            shape_type = shape.shape_type

            if 'PLACEHOLDER' in str(shape_type) and (need_title or need_content):
                shape = Placeholder(shape)
            elif 'PICTURE' in str(shape_type) and need_picture:
                shape = Picture(shape,picture_idx)
                picture_idx += 1
            elif 'CHART' in str(shape_type) and need_chart:
                shape = Chart(shape)
            elif 'TABLE' in str(shape_type) and need_table:
                shape = Table(shape)
            elif 'TEXT_BOX' in str(shape_type) and (need_textbox or (global_args.dataset=='long' and (need_title or need_content))):
                shape = Textbox(shape,textbox_idx)
                textbox_idx += 1
            elif 'AUTO_SHAPE' in str(shape_type) and need_shape:
                shape = AutoShape(shape)
            else:
                continue

            s += shape.discription
            if need_position:
                s += shape.size_info
            if need_text and not (shape.text_info is None):
                s += shape.text_info
            if need_style and not (shape.style_info is None):
                s += shape.style_info
            if need_position and not (shape.space_info is None):
                s += shape.space_info
            s += '\n'
    return s

def get_content_by_instructions(ppt_path, instruction, args, ppt):
    global slides
    global global_args
    prompt = ""
    if not ppt is None:
        slides = ppt.slides
    else:
        ppt = Presentation(ppt_path)
        slides = ppt.slides
    global_args = args

    s = f"There are {len(ppt.slides)} slides with slide height {ppt.slide_height//SCALE} and slide width {ppt.slide_width//SCALE}.\n"

    if args.content_selection:
        try:
            prompt = PPT_content_selection_prompt.format(instruction)
            ppt_operation = query_azure_openai(prompt, model=args.model).strip()
            print('#### Content Selection:')
            print(ppt_operation)
            ppt_content = eval(ppt_operation)
        except:
            print('Content Selection Failed!')
            ppt_operation = "get_content(need_text=1,need_style=1,need_position=1,need_title=1,need_content=1,need_picture=1,need_table=1,need_chart=1,need_textbox=1,need_shape=1)"
            ppt_content = eval(ppt_operation)
    else:
        ppt_operation = "get_content(need_text=1,need_style=1,need_position=1,need_title=1,need_content=1,need_picture=1,need_table=1,need_chart=1,need_textbox=1,need_shape=1)"
        ppt_content = eval(ppt_operation)
    s += ppt_content
    return s, prompt


def eval_get_contents(need_text=True, need_style=True, need_position=True, need_shape_list=None, ppt=None):
    slides = ppt.slides
    s = ""

    idx = 0
    for slide in slides:
        s += f"Slide {idx} with background color {get_fill_color(slide.background)}:\n"
        if slide.notes_slide.notes_text_frame.text:
            s += f"Notes: {slide.notes_slide.notes_text_frame.text}\n"
        idx += 1
        textbox_idx = 0
        picture_idx = 0
        for shape in slide.shapes:
            if need_shape_list is not None and not hasshape(str(shape.shape_type), need_shape_list):
                continue
            if 'PLACEHOLDER' in str(shape.shape_type):
                shape = Placeholder(shape)
            elif 'PICTURE' in str(shape.shape_type):
                shape = Picture(shape,picture_idx)
                picture_idx += 1
            elif 'CHART' in str(shape.shape_type):
                shape = Chart(shape)
            elif 'TABLE' in str(shape.shape_type):
                shape = Table(shape)
            elif 'TEXT_BOX' in str(shape.shape_type):
                shape = Textbox(shape,textbox_idx)
                textbox_idx += 1
            elif 'AUTO_SHAPE' in str(shape.shape_type):
                shape = AutoShape(shape)
            else:
                continue
            s += shape.discription
            if need_position:
                s += shape.size_info
            if need_text and not (shape.text_info is None):
                s += shape.text_info
            if need_style and not (shape.style_info is None):
                s += shape.style_info
            if need_position and not (shape.space_info is None):
                s += shape.space_info
            s += '\n'
    return s

# def get_content_by_instructions(ppt_path, instruction, args, ppt):
#     global slides
#     prompt = ""
#     if not ppt is None:
#         slides = ppt.slides
#     else:
#         ppt = Presentation(ppt_path)
#         slides = ppt.slides

#     s = f"There are {len(ppt.slides)} slides with slide height {ppt.slide_height//SCALE} and slide width {ppt.slide_width//SCALE}.\n"

#     debug_string = ""
#     # select shape
#     shape_list = ['PLACEHOLDER', 'PLACEHOLDER', 'PICTURE', 'TABLE', 'CHART', 'TEXT_BOX', 'AUTO_SHAPE']
#     if args.shape_selection:
#         selected_shape_type = select_shape_type(instruction,args)
#         need_shape_list = [shape_list[i] for i in range(len(selected_shape_type)) if selected_shape_type[i] == 1]
#         if args.dataset == 'long':
#             if 'PLACEHOLDER' in need_shape_list:
#                 need_shape_list.remove('PLACEHOLDER')
#                 need_shape_list.append('TEXT_BOX')
#         print('need shape')
#         print(need_shape_list)
#         debug_string += f"Need Shape: {need_shape_list}\n"
#     else:
#         need_shape_list = shape_list
#     # print(f"== Need Shapes ==\n{need_shape_list}\n\n")

#     # select information
#     if args.info_selection:
#         need_text, need_style, need_position = select_information_type(instruction,args)
#         print(f'need info: text:{need_text},style:{need_style},position:{need_position}')
#         debug_string += f"Need Info: text:{need_text}, style:{need_style}, position:{need_position}\n"
#     else:
#         need_text, need_style, need_position = 1,1,1
#     # print(f"== Need Text == {need_text}\n== Need Style == {need_style}\n== Need Position == {need_position}\n")
    

#     ppt_content = get_contents(need_text=need_text, need_style=need_style, need_position=need_position, need_shape_list=need_shape_list, ppt=ppt)
#     s += ppt_content
#     return s, prompt, debug_string


if __name__ == "__main__":
    pass

