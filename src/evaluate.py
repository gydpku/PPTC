from src import api_doc
from src import prompt_factor
import mosestokenizer 
from src import ppt_reader, utils
from pptx import Presentation
from src import pptx_check
from sacremoses import MosesTokenizer
import os
from tqdm import tqdm

def calc_token_cost(path):
    text = open(path,'r').read()    
    tokenizer = MosesTokenizer()
    tokens = tokenizer.tokenize(text)
    return len(tokens)

def calc_acc(label_path, pred_path, instruction, additional_restrictions=[]):
    pos_total, pos_correct, str_correct = 0,0,0
    # position
    splitted = instruction.split('##')
    instruction, restrictions = splitted[0], splitted[1:]
    restrictions += additional_restrictions
    if len(restrictions) > 0:
        pos_total = 1
        pos_correct = 1
    # try:
    ppt = Presentation(pred_path)
    for res in restrictions:
        slide_id,A,B,rel = [x.strip(" ") for x in res.split(",")]
        try:
            slide = ppt.slides[int(slide_id)]
            pos_correct *= pptx_check.check(slide, A, B, rel)
        except:
            print(res)
            print(instruction)
            pos_correct = 0
        
    
    # string
    label_string = ppt_reader.eval_get_contents(need_text=True, need_style=True, need_position=False,need_shape_list=None,ppt=Presentation(label_path))
    pred_string = ppt_reader.eval_get_contents(need_text=True, need_style=True, need_position=False,need_shape_list=None,ppt=Presentation(pred_path))
    if label_string == pred_string:
        if len(restrictions) > 0:
            if pos_correct == 1:
                str_correct = 1
        else:
            str_correct = 1
    
    print(f'String correct : {str_correct}')
    # except:
    #     print(f'PPT Loading Failed!!!!!')
    #     str_correct, pos_correct = 0,0

    
    return str_correct, pos_total, pos_correct

def check_eval(args):
    token_costs, api_costs = [],[]
    string_total, string_correct, position_total, position_correct = 0,0,0,0
    if args.tf:
        ans = []
        for idx in sorted(os.listdir(args.save_path)):
            if not os.path.isdir(os.path.join(args.save_path,str(idx))):
                continue
            for step in sorted(os.listdir(os.path.join(args.save_path, idx))):
                if '.' in step or not os.path.isdir(os.path.join(args.save_path,str(idx),str(step))):
                    continue
                current_path = os.path.join(args.save_path,str(idx),str(step))
                pred_path = os.path.join(current_path, f"{args.exp_name}_after_pred.pptx")
                label_path = os.path.join(current_path, "after_label.pptx")
                instruction_path = os.path.join(current_path, "instruction.txt")
                prompt_path = os.path.join(current_path, f"{args.exp_name}_prompt.txt")
                api_pred_path = os.path.join(current_path, f"{args.exp_name}_api_pred.txt")

                if calc_api_cost(api_pred_path,args) == 0:
                    ans.append(f'{idx}/{step}')
        print(ans)

def get_error_case(args):
    f = open(f'{args.dataset}_{args.exp_name}_error_info.txt','w+')
    if args.tf:
        for idx in sorted(os.listdir(args.save_path)):
            if not os.path.isdir(os.path.join(args.save_path,str(idx))):
                continue
            for step in sorted(os.listdir(os.path.join(args.save_path, idx))):
                if '.' in step or not os.path.isdir(os.path.join(args.save_path,str(idx),str(step))):
                    continue
                current_path = os.path.join(args.save_path,str(idx),str(step))
                pred_path = os.path.join(current_path, f"{args.exp_name}_after_pred.pptx")
                label_path = os.path.join(current_path, "after_label.pptx")
                instruction_path = os.path.join(current_path, "instruction.txt")
                prompt_path = os.path.join(current_path, f"{args.exp_name}_prompt.txt")
                api_pred_path = os.path.join(current_path, f"{args.exp_name}_api_pred.txt")
                api_label_path = os.path.join(current_path, f"api_label.txt")

                str_c, pos_t, pos_c = calc_acc(label_path, pred_path, instruction_path)
                instruction = open(instruction_path,'r').read().strip()
                labels = open(api_label_path,'r').read().strip()
                preds = open(api_pred_path,'r').read().strip()
                if str_c != 1:
                    f.write(f"{idx}/{step}\n")
                    f.write(instruction+'\n')
                    f.write('## Label ##\n')
                    f.write(labels+'\n')
                    f.write('## Pred ##\n')
                    f.write(preds+'\n\n')
                    f.flush()

def eval(args):
    set_name = 'Create_new_slides' if args.dataset == 'short' else 'Edit_PPT_template'
    token_costs, api_costs = [],[]
    string_total, string_correct, position_total, position_correct, string_acc, position_acc = 0,0,0,0,0,0

    if args.tf:
       # turn_num=0
        turn_nums=[]
        for sess_id, session_path in enumerate(utils.sorted_list(args.user_path+f'PPT_test_output/{set_name}')):
            #print(len(list(utils.sorted_list(args.user_path+f'PPT_test_output/{set_name}'))))
            if not session_path.startswith(args.exp_name):
                continue
            session = utils.parse_test_json(args.user_path+f'PPT_test_output/{set_name}/{session_path}')
            turn_num = 0
            for turn_id, turn in tqdm(enumerate(session)):
                turn_num+=1

                turn_id, instruction, label_api, reply, pred_api, pred_ppt_path, label_ppt_path, prompt_path = turn
                api_costs.append(len(pred_api))
                token_costs.append(calc_token_cost(args.user_path+prompt_path))
                str_c, pos_t, pos_c = calc_acc(args.user_path+label_ppt_path, args.user_path+pred_ppt_path, instruction)
                string_total += 1
                string_correct += str_c
                position_total += pos_t
                position_correct += pos_c
            turn_nums.append(turn_num)
            print(sess_id, turn_nums)
        
        avg_api_costs = sum(api_costs) / len(api_costs)
        avg_token_costs = sum(token_costs) / len(token_costs)
        string_acc = string_correct / string_total
        position_acc = position_correct / position_total

        print(f"avg api cost: {avg_api_costs}")
        print(f"avg_token_costs: {avg_token_costs}")
        print(f"string acc: {string_correct}/{string_total}={string_acc}")
        print(f"position acc: {position_correct}/{position_total}={position_acc}")

    elif args.sess:
        for sess_id, session_path in enumerate(utils.sorted_list(args.user_path+f'PPT_test_output/{set_name}')):
            if not session_path.startswith(args.exp_name):
                continue
            session = utils.parse_test_json(args.user_path+f'PPT_test_output/{set_name}/{session_path}')
            restrictions = []
            for turn_id, turn in tqdm(enumerate(session)):
                turn_id, instruction, label_api, reply, pred_api, pred_ppt_path, label_ppt_path, prompt_path = turn
                splitted = instruction.split('##')
                restrictions.extend(splitted[1:])
                api_costs.append(len(pred_api))
                token_costs.append(calc_token_cost(args.user_path+prompt_path))
                print(turn_id,len(session)-1)
                if turn_id == len(session)-1:
                    str_c, pos_t, pos_c = calc_acc(args.user_path+label_ppt_path, args.user_path+pred_ppt_path, instruction, restrictions)

                    string_total += 1
                    print("add!xxxxxxxxxxxxxxxxxxxxxxxxx", string_total)
                    string_correct += str_c
                    position_total += pos_t
                    position_correct += pos_c
        print(string_total,'string_total')
        try:
            avg_api_costs = sum(api_costs) / string_total
            avg_token_costs = sum(token_costs) / string_total
            string_acc = string_correct / string_total
            position_acc = position_correct / position_total
        except:
            avg_api_costs = sum(api_costs)
            avg_token_costs = sum(token_costs)
            string_acc = string_correct
            position_acc = position_correct 

        print(f"avg api cost: {avg_api_costs}")
        print(f"avg_token_costs: {avg_token_costs}")
        print(f"string acc: {string_correct}/{string_total}={string_acc}")
        print(f"position acc: {position_correct}/{position_total}={position_acc}")
