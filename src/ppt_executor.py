import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from src import api_doc

SLIDE_HEIGHT = 6858000
SLIDE_WIDTH = 9144000
CENTER_TOP = 3429000
CENTER_LEFT = 4572000
SHAPE_HEIGHT = 900000
SHAPE_WIDTH = 900000
TABLE_HEIGHT = 370000 # per line
CONTENT_HEIGHT = 4351338
CONTENT_WIDTH = 7886700
CONTENT_LEFT = 628650
CONTENT_TOP = 1825625
TITLE_HEIGHT = 1325563
TITLE_WIDTH = 7886700
TITLE_LEFT = 628650
TITLE_TOP = 365126
MARGIN = 600000
CORNER_LEFT = 0 + MARGIN
CORNER_TOP = 0 + MARGIN
CORNER_RIGHT = SLIDE_WIDTH - MARGIN
CORNER_BOTTOM = SLIDE_HEIGHT - MARGIN
SHAPE_LEFT = CENTER_LEFT - SHAPE_WIDTH / 2
SHAPE_TOP = CENTER_TOP - SHAPE_HEIGHT / 2

PIC_LEFT = CONTENT_LEFT
PIC_TOP = CONTENT_TOP 

PIC_PATH = "./PPTC/"+"test/pics"

current_shape = None
current_slide = None
slides = None
table = None
picture = None
chart = None
smartart = None
textbox = None
shape = {}
prs = None

color2hex = {
    "blue": "0000FF","light blue": "ADD8E6","dark blue": "00008B",
    "green": "008000","light green": "90EE90","dark green": "006400",
    "yellow": "FFFF00","light yellow": "FFFFE0","dark yellow": "BDB76B",
    "orange": "FFA500","light orange": "FFDAB9","dark orange": "FF8C00",
    "red": "FF0000","light red": "FFC0CB","dark red": "8B0000",
    "black": "000000","white": "FFFFFF","purple": "800080","pink": "FFC0CB",
}

def check_api_in_list(line, api_list):
    for api in api_list:
        if api in line:
            return 1
    return 0

# apis
def API_executor(lines, test=False,args=None):
    # print(f"EXECUTING API:{lines}")
    error_info = ""
    for line in lines:
        if not test:
            if check_api_in_list(line, ["set_left","set_top","set_right","set_bottom"]):
                continue
            if args.api_lack:
                print('Lacking API Exec!')
                if not check_api_in_list(line, api_doc.original_apis):
                    continue
        try:
            if args.dataset == 'long' and line == 'choose_title()':
                eval("choose_textbox(0)")
            if args.api_lack:
                if line == 'seek_assistance()':
                    continue
                elif not check_api_in_list(line, api_doc.original_apis):
                    eval("move_to_slide(0)")
                    eval("insert_note('@@@@@@@@@@')")
                else:
                    eval(line)
            elif args.api_update:
                if check_api_in_list(line, api_doc.update_apis):
                    eval("move_to_slide(0)")
                    eval("insert_note('@@@@@@@@@@')")
                else:
                    eval(line)
            else:
                eval(line) 
        except Exception as e:
            print(line)
            print(f"ERROR: {e}")
            error_info += f"ERROR: {e}\n"
    return error_info
    

def set_ppt(ppt_path):
    global slides, prs, current_shape, current_slide, table, picture, chart, smartart, textbox, shape
    prs = Presentation(ppt_path)
    current_shape = None
    current_slide = None
    slides = None
    table = None
    picture = None
    chart = None
    smartart = None
    textbox = None
    shape = {}
    slides = prs.slides

def set_current_slide(idx):
    global current_slide, slides
    try:
        current_slide = slides[idx]
    except:
        pass

def get_ppt():
    global prs
    return prs

def save_ppt(ppt_path):
    global slides, prs
    prs.save(ppt_path)

def get_current_page_id():
    global current_slide
    return slides.index(current_slide)

# slide
def create_slide():
    global current_slide, current_shape
    slide = slides.add_slide(prs.slide_layouts[1])
    current_slide = slide
    current_shape = None 

def move_to_next_slide():
    global current_slide, current_shape
    current_slide_id = current_slide.slide_id
    try:
        slide = slides[current_slide_id + 1]
    except:
        slide = None
    if slide != None:
        current_slide = slide
        current_shape = None

def move_to_previous_slide():
    global current_slide, current_shape
    current_slide_id = current_slide.slide_id
    try:
        slide = slides[current_slide_id - 1]
    except:
        slide = None
    if slide != None:
        current_slide = slide
        current_shape = None

def move_to_slide(idx):
    global current_slide, current_shape
    try:
        slide = slides[idx]
    except:
        slide = None
    if slide != None:
        current_slide = slide
        current_shape = None

def set_background_color(color):
    global current_slide, current_shape
    current_slide.background.fill.solid()
    current_slide.background.fill.fore_color.rgb = RGBColor.from_string(color2hex[color])

# choose
def choose_title():
    global current_slide, current_shape
    current_shape = current_slide.shapes.title

def choose_content():
    global current_slide, current_shape
    current_shape = current_slide.placeholders[1]

def choose_textbox(idx=0):
    global current_slide, current_shape, textbox
    if textbox != None:
        current_shape = textbox
    else:
        cur_idx = 0
        for shape in current_slide.shapes:
            if 'TEXT_BOX' in str(shape.shape_type):
                if cur_idx == idx:
                    current_shape = shape
                    break
                else:
                    cur_idx += 1

def choose_picture(idx=0):
    global current_slide, current_shape, picture
    if picture != None:
        current_shape = picture
    else:
        cur_idx = 0
        for shape in current_slide.shapes:
            if 'PICTURE' in str(shape.shape_type):
                if cur_idx == idx:
                    current_shape = shape
                    break
                else:
                    cur_idx += 1

def choose_chart():
    global current_slide, current_shape, chart
    if chart != None:
        current_shape = chart
    else:
        for shape in current_slide.shapes:
            if 'CHART' in str(shape.shape_type):
                current_shape = shape

def choose_shape(shape_name):
    global current_slide, current_shape, shape
    if shape[shape_name] != None:
        current_shape = shape[shape_name]
    else:
        for shape in current_slide.shapes:
            if 'AUTO_SHAPE' in str(shape.shape_type):
                current_shape = shape

def choose_table():
    global current_slide, current_shape, table
    if table != None:
        current_shape = table
    else:
        for shape in current_slide.shapes:
            if 'TABLE' in str(shape.shape_type):
                current_shape = shape
                table = current_shape.table

def choose_table_cell(row_id, column_id):
    global current_slide, current_shape, table
    cell = table.cell(row_id, column_id)
    current_shape = cell

# text
def insert_text(text):
    global current_slide, current_shape
    try:
        current_shape.text_frame.text += text
    except:
        try:
            current_shape.text += text
        except:
            pass
    
    return

def insert_bullet_point(text):
    global current_slide, current_shape
    try:
        p = current_shape.text_frame.add_paragraph()
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        p.text = text
        p.level = 0
    except:
        pass

def insert_note(note):
    global current_slide, current_shape
    current_slide.notes_slide.notes_text_frame.text += note

def insert_textbox():
    global current_slide, current_shape, textbox
    textbox = current_slide.shapes.add_textbox(CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, TABLE_HEIGHT)
    current_shape = textbox

def delete_text():
    global current_slide, current_shape
    try:
        current_shape.text_frame.text = "" 
    except:
        try:
            current_shape.text = ""
        except:
            pass

def set_font_size(size):
    global current_slide, current_shape
    for paragraph in current_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size)

def set_font_color(color):
    global current_slide, current_shape
    for paragraph in current_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor.from_string(color2hex[color])

def set_font_bold():
    global current_slide, current_shape
    for paragraph in current_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.bold = True

def set_font_italic():
    global current_slide, current_shape
    for paragraph in current_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.italic = True

def set_font_underline():
    global current_slide, current_shape
    for paragraph in current_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.underline = True

def set_font_style(font_name):
    global current_slide, current_shape
    for paragraph in current_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name

def set_line_space(line_space_level=0):
    global current_slide, current_shape
    for paragraph in current_shape.text_frame.paragraphs:
        paragraph.line_spacing = line_space_level

def text_align_left():
    global current_slide, current_shape
    for paragraph in current_shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT

def text_align_center():
    global current_slide, current_shape
    for paragraph in current_shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER

def text_align_right():
    global current_slide, current_shape
    for paragraph in current_shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.RIGHT

# shape
def insert_rectangle():
    global current_slide, current_shape, shape
    current_shape = current_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, SHAPE_LEFT, SHAPE_TOP, SHAPE_WIDTH, SHAPE_HEIGHT)
    shape[f'rectangle'] = current_shape

def insert_right_arrow():
    global current_slide, current_shape, shape
    current_shape = current_slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, SHAPE_LEFT, SHAPE_TOP, SHAPE_WIDTH, SHAPE_HEIGHT)
    shape[f'right_arrow'] = current_shape

def insert_rounded_rectangle():
    global current_slide, current_shape, shape
    current_shape = current_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, SHAPE_LEFT, SHAPE_TOP, SHAPE_WIDTH, SHAPE_HEIGHT)
    shape[f'rounded_rectangle'] = current_shape

def insert_triangle():
    global current_slide, current_shape, shape
    current_shape = current_slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, SHAPE_LEFT, SHAPE_TOP, SHAPE_WIDTH, SHAPE_HEIGHT)
    shape[f'triangle'] = current_shape

def insert_callout():
    global current_slide, current_shape, shape
    current_shape = current_slide.shapes.add_shape(MSO_SHAPE.BALLOON, SHAPE_LEFT, SHAPE_TOP, SHAPE_WIDTH, SHAPE_HEIGHT)
    shape[f'callout'] = current_shape

def insert_cloud():
    global current_slide, current_shape, shape
    current_shape = current_slide.shapes.add_shape(MSO_SHAPE.CLOUD, SHAPE_LEFT, SHAPE_TOP, SHAPE_WIDTH, SHAPE_HEIGHT)
    shape[f'cloud'] = current_shape

def insert_star():
    global current_slide, current_shape, shape
    current_shape = current_slide.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, SHAPE_LEFT, SHAPE_TOP, SHAPE_WIDTH, SHAPE_HEIGHT)
    shape[f'star'] = current_shape

def insert_circle():
    global current_slide, current_shape, shape
    current_shape = current_slide.shapes.add_shape(MSO_SHAPE.OVAL, SHAPE_LEFT, SHAPE_TOP, SHAPE_WIDTH, SHAPE_HEIGHT)
    shape[f'circle'] = current_shape

# picture
def insert_picture(picture_name):
    global current_slide, current_shape, picture
    try:
        current_shape = current_slide.shapes.add_picture(f"{PIC_PATH}/{picture_name}.png", PIC_LEFT, PIC_TOP)
    except:
        current_shape = current_slide.shapes.add_picture(f"{PIC_PATH}/none.png", PIC_LEFT, PIC_TOP)
    picture = current_shape
    picture.width = Cm(6)
    picture.height = Cm(6)

# basic
def set_width(width):
    global current_slide, current_shape
    current_shape.width = int(width * 360000)

def set_height(height):
    global current_slide, current_shape
    current_shape.height = int(height * 360000)

def rotate_element(angle):
    global current_slide, current_shape
    current_shape.rotation = angle

def set_fill_color(color):
    global current_slide, current_shape
    current_shape.fill.solid()
    current_shape.fill.fore_color.rgb = RGBColor.from_string(color2hex[color])

def align_top_right_corner():
    global current_slide, current_shape
    current_shape.left = CORNER_RIGHT - current_shape.width
    current_shape.top = CORNER_TOP

def align_top_left_corner():
    global current_slide, current_shape
    current_shape.left = CORNER_LEFT
    current_shape.top = CORNER_TOP

def align_bottom_right_corner():
    global current_slide, current_shape
    current_shape.left = CORNER_RIGHT - current_shape.width
    current_shape.top = CORNER_BOTTOM - current_shape.height

def align_bottom_left_corner():
    global current_slide, current_shape
    current_shape.left = CORNER_LEFT
    current_shape.top = CORNER_BOTTOM - current_shape.height

def align_slide_left():
    global current_slide, current_shape
    current_shape.left = CORNER_LEFT

def align_slide_right():
    global current_slide, current_shape
    current_shape.left = CENTER_LEFT + MARGIN

def align_slide_top():
    global current_slide, current_shape
    current_shape.top = CORNER_TOP

def align_slide_bottom():
    global current_slide, current_shape
    current_shape.top = CENTER_TOP + MARGIN

def align_slide_center():
    global current_slide, current_shape
    current_shape.top = CENTER_TOP - current_shape.height // 2
    current_shape.left = CENTER_LEFT - current_shape.width // 2

def set_left(left):
    global current_slide, current_shape
    current_shape.left = left

def set_top(top):
    global current_slide, current_shape
    current_shape.top = top

# table
def insert_table(row_num, col_num):
    global current_slide, current_shape, table, TABLE_HEIGHT, TABLE_WIDTH
    table_height = row_num * TABLE_HEIGHT
    current_shape = current_slide.shapes.add_table(row_num, col_num, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, table_height).table
    table = current_shape

def insert_table_row(row_data):
    global current_slide, current_shape
    row_num = len(current_shape.rows)
    col_num = len(current_shape.columns)
    idx = 0
    for i in range(row_num):
        if table.cell(i, 0).text == "":
            idx = i
            break

    for j in range(len(row_data)):
        table.cell(idx, j).text = row_data[j]

# chart
def insert_line_chart(data,series=None):
    global current_slide, current_shape, chart
    chart_data = ChartData()
    if series is None:
        series = [f'Series {i}' for i in range(len(data))]
    for i in range(len(data)):
        chart_data.add_series(series[i], [data[i]])
    chart_data.categories = [f'Category {i}' for i in range(len(data))]
    current_shape = current_slide.shapes.add_chart(XL_CHART_TYPE.LINE, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT, chart_data).chart
    chart = current_shape
    return chart

def insert_bar_chart(data,series=None):
    global current_slide, current_shape, chart
    chart_data = ChartData()
    if series is None:
        series = [f'Series {i}' for i in range(len(data))]
    for i in range(len(data)):
        chart_data.add_series(series[i], [data[i]])
    chart_data.categories = [f'Category {i}' for i in range(len(data))]
    current_shape = current_slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT, chart_data).chart
    chart = current_shape
    return chart

def insert_pie_chart(data,series=None):
    global current_slide, current_shape, chart
    chart_data = ChartData()
    if series is None:
        series = [f'Series {i}' for i in range(len(data))]
    for i in range(len(data)):
        chart_data.add_series(series[i], [data[i]])
    chart_data.categories = [f'Category {i}' for i in range(len(data))]
    current_shape = current_slide.shapes.add_chart(XL_CHART_TYPE.PIE, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT, chart_data).chart
    chart = current_shape
    return chart

def set_chart_title(title):
    global current_slide, current_shape
    current_shape.chart_title.text_frame.text = title

if __name__ == '__main__':
    prs = Presentation()
    slides = prs.slides
    # slide
    create_slide()
    insert_table(2,3)
    insert_table_row(["hi1","hi2","hi3"])
    # save
    prs.save('test_table1.pptx')