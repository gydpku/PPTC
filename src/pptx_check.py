import sys
import os
import json
from src import ppt_reader
from pptx import Presentation
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from tqdm import tqdm

SLIDE_HEIGHT_MIDDLE = 3429000
SLIDE_WIDTH_MIDDLE = 4572000

def check_left(A,B):
    return int(A.left+A.width < B.left)
def check_top(A,B):
    return int(A.top+A.height < B.top)
def check_right(A,B):
    return int(A.left > B.left+B.width)
def check_bottom(A,B):
    return int(A.top > B.top+B.height)
def check_slide_left(A):
    return int(A.left+A.width < SLIDE_WIDTH_MIDDLE)
def check_slide_top(A):
    return int(A.top+A.height < SLIDE_HEIGHT_MIDDLE)
def check_slide_right(A):
    return int(A.left > SLIDE_WIDTH_MIDDLE)
def check_slide_bottom(A):
    return int(A.top > SLIDE_HEIGHT_MIDDLE)
def check_slide_center(A):
    return int(A.left < SLIDE_WIDTH_MIDDLE and A.top < SLIDE_HEIGHT_MIDDLE and A.left+A.width > SLIDE_WIDTH_MIDDLE and A.top+A.height > SLIDE_HEIGHT_MIDDLE)
def choose_table(slide):
    for shape in slide.shapes:
        if 'TABLE' in str(shape.shape_type):
            return shape
    return None
def choose_chart(slide):
    for shape in slide.shapes:
        if 'CHART' in str(shape.shape_type):
            return shape
    return None  
def choose_content(slide):
    for shape in slide.shapes:
        if "Content" in shape.name:
            return shape
    return None
def choose_title(slide):
    for shape in slide.shapes:
        if "Title" in shape.name:
            return shape
    return None
def choose_shape(slide, shape_name):
    shape_name_dict = {
        "rectangle": "Rectangle",
        "right_arrow": "Right Arrow",
        "arrow": "Right Arrow",
        "cloud": "Cloud",
        "rounded_rectangle": "Rounded Rectangle",
        "triangle": "Isosceles",
        "callout": "Callout",
        "star": "Star",
        "circle": "Oval"
    }
    shape_name = shape_name_dict[shape_name]
    for shape in slide.shapes:
        if shape_name in shape.name:
            if shape_name == "Rectangle":
                if not "Rounded Rectangle" in shape.name:
                    return shape 
            else: 
                return shape
    return None
def choose_picture(slide,idx=0):
    cur_idx = 0
    for shape in slide.shapes:
        if 'PICTURE' in str(shape.shape_type):
            if cur_idx == idx:
                return shape
            else:
                cur_idx += 1
    return None
def choose_textbox(slide,idx=0):
    cur_idx = 0
    for shape in slide.shapes:
        if 'TEXT_BOX' in str(shape.shape_type):
            if cur_idx == idx:
                return shape
            else:
                cur_idx += 1
    return None

def choose_object(slide, object_name):
    if object_name == 'table':
        return choose_table(slide)
    elif 'chart' in object_name:
        return choose_chart(slide)
    elif 'textbox' in object_name:
        object_id = int(object_name[8:]) if object_name!='textbox' else 0
        return choose_textbox(slide,object_id)
    elif object_name == 'content':
        return choose_content(slide)
    elif object_name == 'title':
        return choose_title(slide)
    elif ('picture' in object_name):
        object_id = int(object_name[8:]) if object_name!='picture' else 0
        return choose_picture(slide,object_id)
    elif object_name == 'image':
        return choose_picture(slide)
    elif object_name in ['rectangle', 'right_arrow', 'rounded_rectangle', 'triangle', 'callout', 'star', 'circle', 'arrow', 'cloud']:
        return choose_shape(slide, object_name)
    else:
        return None

def check(slide, A, B, rel):
    result = 1
    if B == "slide":
        A = choose_object(slide, A)
        if A is None:
            print("No object found")
            return 0
        if "left" in rel:
            result *= check_slide_left(A)
        if "top" in rel:
            result *= check_slide_top(A)
        if "right" in rel:
            result *= check_slide_right(A)
        if "bottom" in rel:
            result *= check_slide_bottom(A)
        if "center" in rel:
            result *= check_slide_center(A)
    else:
        A = choose_object(slide, A)
        B = choose_object(slide, B)
        if A is None or B is None:
            print("No object found")
            return 0
        if "left" in rel:
            result *= check_left(A,B)
        if "top" in rel:
            result *= check_top(A,B)
        if "right" in rel:
            result *= check_right(A,B)
        if "bottom" in rel:
            result *= check_bottom(A,B)
    return result